import datetime
import os
import json
from io import BytesIO

import markdown
from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.text.text import Font

from core.config import settings
from mdtree.parser import parse_string, Out, Heading
from mdtree.utils import get_random_theme, get_random_file


class Tree2PPT:
    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None

    def __init__(self, markdown_str):
        self.init_pptx()
        self.init_markdown(markdown_str)
        self.traverse_tree(self.tree)

        now = datetime.datetime.now().timestamp()
        file_name = "ppt_" + str(now).replace(".", "") + ".pptx"
        file_path = os.path.join(settings.PPT_PATH, file_name)
        self.current_path = file_path
        self.prs.save(file_path)

    def init_pptx(self):
        prs = Presentation()
        self.theme = get_random_theme()
        self.prs = prs

    def init_markdown(self, md_str):
        self.md_str = md_str
        self.out = parse_string(md_str)
        self.tree = self.out.main

    def traverse_tree(self, heading):
        if heading is not None and (heading.source is None or heading.source == ""):
            content = ""
            if heading.children != []:
                for child in heading.children:
                    content = content + child.text + "\n"
            MD2Slide(self.prs, self.theme, heading.text, content=content)
        elif heading is not None:
            MD2Slide(self.prs, self.theme, heading.text, content=heading.source)
        else:
            return

        # self.make_slide_demo(self.prs, heading.text, heading.source)
        if heading.children != []:
            for child in heading.children:
                self.traverse_tree(child)

    def get_file_path(self):
        return self.current_path

    def save_stream(self):
        stream = BytesIO()
        self.prs.save(stream)
        stream.seek(0)  # Reset the stream position to the beginning
        return stream

    def extract_shape_properties(self, shape):
        """Extract text properties from a shape."""
        properties = {}
        if shape.has_text_frame:
            tf = shape.text_frame
            properties["text"] = shape.text.strip()
            properties["alignment"] = {
                "horizontal": tf.paragraphs[0].alignment.name
                if tf.paragraphs[0].alignment
                else None,
                "vertical": tf.vertical_anchor if tf.vertical_anchor else None,
            }
            properties["runs"] = []

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    run_properties = {
                        "text": run.text,
                        "font": font.name,
                        "size": font.size.pt if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                        "underline": font.underline,
                        "color": self.get_font_color(font),
                    }
                    properties["runs"].append(run_properties)
        return properties

    def get_font_color(self, font):
        """Safely get the RGB color or theme color of the font."""
        if font.color:
            if hasattr(font.color, "rgb") and font.color.rgb:
                # If the font color has an RGB property
                return f"#{font.color.rgb:06X}"
            elif hasattr(font.color, "theme_color") and font.color.theme_color:
                # If the font color is a theme color
                return f"Theme Color: {font.color.theme_color}"
        # If no explicit color is set
        return "Default Color"

    def to_markdown(self):
        """Convert presentation content to Markdown format with text properties."""
        md_content = []
        for idx, slide in enumerate(self.prs.slides):
            md_content.append(f"# Slide {idx + 1}\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_props = self.extract_shape_properties(shape)
                    for run in shape_props.get("runs", []):
                        color_hex = run["color"] if run["color"] else "default"
                        md_content.append(
                            f"- **Text**: {run['text']} | "
                            f"**Font**: {run['font']} | **Size**: {run['size']}pt | "
                            f"**Color**: {color_hex} | **Bold**: {run['bold']} | "
                            f"**Italic**: {run['italic']} | **Underline**: {run['underline']}"
                        )
            md_content.append("")  # Add blank line between slides
        return "\n".join(md_content)

    def to_json(self):
        """Convert presentation content to JSON format with text properties."""
        slides_content = []
        for idx, slide in enumerate(self.prs.slides):
            slide_data = {"slide_number": idx + 1, "content": []}
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_props = self.extract_shape_properties(shape)
                    slide_data["content"].append(shape_props)
            slides_content.append(slide_data)
        return json.dumps(slides_content, indent=4)


class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"


class MD2Slide:
    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "Arial"
    font_title_size: Pt = Pt(25)
    font_content_size: Pt = Pt(16)
    font_title_color: rgb = RGBColor(51, 0, 102)
    font_content_color: rgb = RGBColor(51, 0, 102)

    def __init__(self, presentation, theme_path, title, content, *args, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[8])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()

    def init_slide(self):
        placeholder1 = self.slide.placeholders[1]
        path = get_random_file(self.theme)
        picture = placeholder1.insert_picture(path)
        placeholder2 = self.slide.placeholders[2]
        placeholder2.element.getparent().remove(placeholder2.element)
        # 2、设置占位符宽高
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        if "font_name" in kwargs:
            self.font_name = kwargs["font_name"]
        if "font_title_size" in kwargs:
            self.font_title_size = kwargs["font_title_size"]
        if "font_content_size" in kwargs:
            self.font_content_size = kwargs["font_content_size"]
        if "font_title_color" in kwargs:
            self.font_title_color = kwargs["font_title_color"]
        if "font_content_color" in kwargs:
            self.font_content_color = kwargs["font_content_color"]

    def get_font(self, font: Font, category: str):
        font.bold = True
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        # Get slide width
        slide_width = self.presentation.slide_width

        # Calculate 80% of the slide width
        title_width = slide_width * 0.9
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(Inches(0.3), Inches(0.3), title_width, Inches(1))
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # 添加标题
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(
            Inches(1), Inches(1.8), Inches(8), Inches(5)
        )
        tf = text_box_content.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        # 添加正文
        paragraph = tf.paragraphs[0]
        paragraph.text = self.content.replace("<p>", "").replace("</p>", "\n")
        self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
        # TODO 处理正文
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def processing_md_str(self, md_str):
        print(md_str)
        md = markdown.Markdown()
        html1 = md.convert(md_str)
        print(html1)
